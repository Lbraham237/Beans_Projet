from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ──────────────────────────────────────────────────────────────────
GREEN_DARK  = RGBColor(0x2E, 0x7D, 0x32)   # vert foncé
GREEN_MID   = RGBColor(0x43, 0xA0, 0x47)   # vert moyen
GREEN_LIGHT = RGBColor(0xC8, 0xE6, 0xC9)   # vert très clair
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x21, 0x21, 0x21)
GREY        = RGBColor(0x75, 0x75, 0x75)
ACCENT      = RGBColor(0xFF, 0xB3, 0x00)   # jaune/or pour highlights

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # layout vide


# ── Helpers ───────────────────────────────────────────────────────────────────
def bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg_ = slide.background
    fill = bg_.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, width, height, text, font_size=18, bold=False,
        color=DARK, bg_color=None, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if bg_color:
        from pptx.oxml.ns import qn
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox

def rect(slide, left, top, width, height, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_box(slide, left, top, width, height, title, bullets,
                   title_color=GREEN_DARK, bullet_color=DARK, bg_color=None, font_size=15):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(font_size + 2)
    run.font.color.rgb = title_color

    for b in bullets:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  {b}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = bullet_color
    return txBox


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)

rect(s, 0, 0, 13.33, 0.08, ACCENT)          # bande haute
rect(s, 0, 7.42, 13.33, 0.08, ACCENT)       # bande basse

box(s, 0.5, 1.2, 12.3, 1.2,
    "Classification de Maladies du Haricot",
    font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.5, 2.5, 12.3, 0.7,
    "Projet ML & Deep Learning — Module PY-ML-DL-M1",
    font_size=20, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

rect(s, 3.5, 3.4, 6.3, 0.05, ACCENT)

box(s, 0.5, 3.7, 12.3, 0.6,
    "MVOGO Abraham  ·  MAAROUFI Abdelhamid",
    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.5, 4.4, 12.3, 0.5,
    "Dataset : Beans — Makerere AI Lab  |  3 classes  |  1 295 images",
    font_size=15, color=GREEN_LIGHT, align=PP_ALIGN.CENTER, italic=True)

# icônes textuels
for i, (emoji, label) in enumerate([("🌿","Angular Leaf Spot"), ("🍂","Bean Rust"), ("✅","Healthy")]):
    x = 2.5 + i * 3.0
    box(s, x, 5.3, 2.6, 0.5, emoji, font_size=28, align=PP_ALIGN.CENTER, color=WHITE)
    box(s, x, 5.85, 2.6, 0.4, label, font_size=13, align=PP_ALIGN.CENTER, color=GREEN_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLÉMATIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Problématique", font_size=24, bold=True, color=WHITE)

box(s, 0.5, 0.9, 12.3, 1.0,
    "Comment diagnostiquer automatiquement la santé d'une feuille de haricot\nà partir d'une simple photo prise au smartphone ?",
    font_size=20, bold=True, color=GREEN_DARK, align=PP_ALIGN.CENTER)

rect(s, 0.5, 2.1, 12.3, 0.04, GREEN_LIGHT)

# 3 cartes
cards = [
    ("🌍 Enjeu agricole", ["Petits exploitants en Afrique de l'Est", "Perte de récoltes par maladies non détectées", "Besoin d'un outil simple et accessible"]),
    ("🎯 Objectif", ["Classifier 3 états : sain, tache angulaire, rouille", "Comparer ML classique vs Deep Learning", "Déployer un dashboard interactif"]),
    ("📊 Approche", ["Baseline ML : Logistic Regression + régularisation", "CNN from scratch", "Transfer Learning MobileNetV2"]),
]
for i, (title, bullets) in enumerate(cards):
    x = 0.4 + i * 4.25
    rect(s, x, 2.4, 3.9, 4.5, GREEN_LIGHT)
    add_bullet_box(s, x+0.15, 2.5, 3.7, 4.3, title, bullets,
                   title_color=GREEN_DARK, bullet_color=DARK, font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Dataset — Beans (Makerere AI Lab, 2020)", font_size=24, bold=True, color=WHITE)

# Stats
stats = [("1 295", "images totales"), ("128×128", "résolution px"), ("3", "classes"), ("HuggingFace", "source")]
for i, (val, label) in enumerate(stats):
    x = 0.5 + i * 3.1
    rect(s, x, 0.8, 2.8, 1.4, GREEN_DARK)
    box(s, x, 0.85, 2.8, 0.7, val, font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    box(s, x, 1.45, 2.8, 0.4, label, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

# Répartition
add_bullet_box(s, 0.5, 2.5, 5.8, 4.5,
    "Répartition des splits",
    ["Train    : 1 034 exemples (79.8%)",
     "Validation :  133 exemples (10.3%)",
     "Test       :  128 exemples (9.9%)",
     "",
     "Distribution équilibrée entre les 3 classes",
     "~344 images par classe"],
    title_color=GREEN_DARK, font_size=15)

add_bullet_box(s, 6.6, 2.5, 6.3, 4.5,
    "Prétraitement appliqué",
    ["Resize automatique → 128×128 pixels",
     "Normalisation [0, 1]",
     "Batching : 32 images/batch",
     "",
     "Augmentation (train uniquement) :",
     "  • RandomFlip horizontal/vertical",
     "  • RandomRotation ±10%",
     "  • RandomZoom ±10%",
     "  • RandomContrast ±10%"],
    title_color=GREEN_DARK, font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE DU PROJET
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Architecture Modulaire du Projet", font_size=24, bold=True, color=WHITE)

modules = [
    ("data_loader.py", "Chargement HuggingFace\nRépartition train/val/test\nVersion NumPy pour ML"),
    ("ml_baseline.py", "LogReg Ridge/Lasso/ElasticNet\nRésultats : accuracy, F1\nDiagnostic biais/variance"),
    ("models.py", "MLP baseline\nCNN from scratch\nTransfer MobileNetV2"),
    ("training.py", "Data augmentation\nCallbacks EarlyStopping\nReduceLROnPlateau"),
    ("evaluation.py", "Matrice de confusion\nCourbes d'apprentissage\nTableau comparatif"),
    ("dashboard.py", "Interface Streamlit\nUpload image\nPrédiction + probabilités"),
]
for i, (name, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.25
    y = 0.9 + row * 3.1
    rect(s, x, y, 3.9, 2.8, GREEN_DARK if row == 0 else GREEN_MID)
    box(s, x+0.1, y+0.1, 3.7, 0.5, name, font_size=15, bold=True, color=ACCENT)
    box(s, x+0.1, y+0.65, 3.7, 2.0, desc, font_size=13, color=WHITE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BASELINE ML
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Baseline ML — Régression Logistique Régularisée", font_size=24, bold=True, color=WHITE)

add_bullet_box(s, 0.5, 0.8, 5.8, 3.5,
    "Modèle & Régularisation",
    ["Logistic Regression (scikit-learn)",
     "",
     "L2 (Ridge) : pénalise les grands poids",
     "  → réduit la variance",
     "L1 (Lasso) : force la parcimonie",
     "  → sélection de features",
     "ElasticNet : combinaison L1 + L2",
     "  → compromis sparsité/stabilité"],
    title_color=GREEN_DARK, font_size=14)

add_bullet_box(s, 6.6, 0.8, 6.2, 3.5,
    "Résultats (LogReg L2, C=1.0)",
    ["Train accuracy  : 100.0%  ⚠️ sur-apprentissage",
     "Val accuracy    :  67.7%",
     "Test accuracy   :  63.3%",
     "F1 macro        :  0.634",
     "",
     "Diagnostic : VARIANCE ÉLEVÉE",
     "Écart train-val = 32.3%",
     "Le modèle mémorise le bruit"],
    title_color=GREEN_DARK, font_size=14)

rect(s, 0.5, 4.5, 12.3, 0.05, GREEN_LIGHT)
box(s, 0.5, 4.65, 12.3, 2.5,
    "Pourquoi la baseline est limitée ?\n\n"
    "Les images 128×128 aplaties → vecteurs de 49 152 features (pixels RGB)\n"
    "La régression logistique ignore la structure spatiale des images\n"
    "→ Le Deep Learning est justifié pour exploiter les patterns visuels locaux",
    font_size=14, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CNN FROM SCRATCH
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Deep Learning — CNN From Scratch", font_size=24, bold=True, color=WHITE)

add_bullet_box(s, 0.5, 0.8, 5.8, 5.8,
    "Architecture CNN",
    ["Input : (128, 128, 3)",
     "",
     "Bloc 1 : Conv2D(32) → BN → ReLU → MaxPool",
     "Bloc 2 : Conv2D(64) → BN → ReLU → MaxPool",
     "Bloc 3 : Conv2D(128) → BN → ReLU → MaxPool",
     "",
     "GlobalAveragePooling2D",
     "Dropout(0.5)",
     "Dense(3, softmax)",
     "",
     "Total : 94 307 paramètres",
     "Optimiseur : Adam (lr=0.001)",
     "Callbacks : EarlyStopping + ReduceLR"],
    title_color=GREEN_DARK, font_size=13)

add_bullet_box(s, 6.6, 0.8, 6.2, 3.0,
    "Résultats",
    ["Test accuracy  : 79.7%",
     "F1 macro       : 0.795",
     "Paramètres     : 94 307",
     "",
     "+16.4 points vs baseline ML"],
    title_color=GREEN_DARK, font_size=14)

add_bullet_box(s, 6.6, 4.0, 6.2, 2.8,
    "Anti Vanishing Gradient",
    ["BatchNormalization après chaque Conv",
     "Activation ReLU (gradient non nul)",
     "GlobalAveragePooling (pas Flatten)",
     "ReduceLROnPlateau si stagnation"],
    title_color=GREEN_DARK, font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRANSFER LEARNING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Deep Learning Avancé — Transfer Learning MobileNetV2", font_size=24, bold=True, color=WHITE)

add_bullet_box(s, 0.5, 0.8, 5.8, 5.8,
    "Architecture Transfer Learning",
    ["Base : MobileNetV2 (pré-entraîné ImageNet)",
     "  → poids figés (trainable=False)",
     "",
     "Tête de classification ajoutée :",
     "  GlobalAveragePooling2D",
     "  Dense(128, ReLU)",
     "  Dropout(0.3)",
     "  Dense(3, softmax)",
     "",
     "Total : 2 261 827 paramètres",
     "Dont 94% non entraînables (base figée)",
     "",
     "Justification : les features visuelles bas",
     "niveau (textures, formes) appris sur",
     "ImageNet sont réutilisables pour les feuilles"],
    title_color=GREEN_DARK, font_size=13)

add_bullet_box(s, 6.6, 0.8, 6.2, 3.0,
    "Résultats",
    ["Test accuracy  : 86.7%",
     "F1 macro       : 0.867",
     "Paramètres     : 2 261 827",
     "",
     "+23.4 points vs baseline ML",
     "+7.0 points vs CNN scratch"],
    title_color=GREEN_DARK, font_size=14)

add_bullet_box(s, 6.6, 4.0, 6.2, 2.8,
    "Pourquoi MobileNetV2 ?",
    ["Architecture légère (mobile-friendly)",
     "Dépthwise separable convolutions",
     "Inverted residuals (bottleneck)",
     "Adapté à des datasets de taille modeste"],
    title_color=GREEN_DARK, font_size=14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COMPARAISON RÉSULTATS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Comparaison des Modèles", font_size=24, bold=True, color=WHITE)

# Tableau
headers = ["Modèle", "Accuracy (test)", "F1 macro", "Paramètres", "Temps"]
rows_data = [
    ["LogReg régularisée (baseline)", "63.3%", "0.634", "36 864", "rapide"],
    ["CNN from scratch",              "79.7%", "0.795", "94 307", "moyen"],
    ["Transfer MobileNetV2",          "86.7%", "0.867", "2 261 827", "moyen"],
]

col_widths = [3.8, 2.2, 1.8, 2.4, 1.5]
col_x = [0.4]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

# headers
for j, (h, x, w) in enumerate(zip(headers, col_x, col_widths)):
    rect(s, x, 0.9, w-0.05, 0.55, GREEN_DARK)
    box(s, x+0.05, 0.95, w-0.1, 0.45, h, font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_data):
    y = 1.55 + i * 0.7
    row_bg = GREEN_LIGHT if i == 2 else (RGBColor(0xF5,0xF5,0xF5) if i % 2 == 0 else WHITE)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_widths)):
        rect(s, x, y, w-0.05, 0.65, row_bg)
        is_best = i == 2 and j in (1, 2)
        box(s, x+0.05, y+0.05, w-0.1, 0.55, ("★ " if is_best else "") + cell,
            font_size=13, bold=is_best, color=GREEN_DARK if is_best else DARK,
            align=PP_ALIGN.CENTER)

# Analyse
box(s, 0.5, 4.7, 12.3, 0.4,
    "Analyse comparative", font_size=16, bold=True, color=GREEN_DARK)
box(s, 0.5, 5.15, 12.3, 2.1,
    "• La baseline ML souffre d'un sur-apprentissage sévère (train=100% vs test=63%) : les pixels bruts ne capturent pas la sémantique visuelle.\n"
    "• Le CNN from scratch gagne +16 pts grâce à l'invariance spatiale des convolutions et la BatchNorm anti-vanishing.\n"
    "• Le Transfer Learning atteint 86.7% en réutilisant les features ImageNet : optimal pour un dataset de 1 295 images.",
    font_size=13, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "Dashboard Streamlit — Déploiement (Jalon 9)", font_size=24, bold=True, color=WHITE)

add_bullet_box(s, 0.5, 0.8, 5.5, 6.3,
    "Fonctionnalités",
    ["Upload d'une image JPG/PNG",
     "Prétraitement automatique :",
     "  • Resize 128×128",
     "  • Normalisation [0,1]",
     "",
     "Affichage de la prédiction",
     "Score de confiance (%)",
     "Détail des 3 probabilités",
     "  (barres de progression)",
     "",
     "Chargement du modèle en cache",
     "  → réponse instantanée",
     "",
     "Lancement : streamlit run app/dashboard.py"],
    title_color=GREEN_DARK, font_size=14)

# Simulation de l'interface
rect(s, 6.3, 0.8, 6.6, 6.3, RGBColor(0xF0, 0xF4, 0xF0))
box(s, 6.4, 0.85, 6.4, 0.45, "🌱 Détection de maladies du haricot",
    font_size=13, bold=True, color=GREEN_DARK)
rect(s, 6.4, 1.4, 6.3, 0.9, RGBColor(0xE0,0xE0,0xE0))
box(s, 6.4, 1.5, 6.3, 0.7, "📁  Choisissez une photo de feuille",
    font_size=12, color=GREY, align=PP_ALIGN.CENTER)
box(s, 6.4, 2.5, 6.3, 0.45, "Prédiction : Tache angulaire 🦠",
    font_size=14, bold=True, color=GREEN_DARK)
box(s, 6.4, 3.05, 3.0, 0.35, "Confiance", font_size=11, color=GREY)
box(s, 6.4, 3.35, 3.0, 0.6, "90.4%", font_size=28, bold=True, color=GREEN_DARK)
box(s, 6.4, 4.1, 6.3, 0.35, "Détail des probabilités :", font_size=12, bold=True, color=DARK)
for idx, (label, pct) in enumerate([("Tache angulaire 🦠", "90.4%"),
                                     ("Rouille 🍂", "6.8%"),
                                     ("Saine ✅", "2.8%")]):
    box(s, 6.4, 4.55 + idx * 0.6, 6.3, 0.35, f"{label} — {pct}", font_size=11, color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CI/CD & REPRODUCTIBILITÉ
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, 13.33, 0.55, GREEN_DARK)
box(s, 0.4, 0.08, 12, 0.4, "CI/CD, Tests & Reproductibilité (Bonus)", font_size=24, bold=True, color=WHITE)

add_bullet_box(s, 0.5, 0.8, 5.8, 3.5,
    "Pipeline CI/CD (.github/workflows/ci.yml)",
    ["Déclenchement : push / pull request",
     "Runner : ubuntu-latest + Python 3.12",
     "",
     "Étapes :",
     "  1. Installation des dépendances légères",
     "     (numpy, pandas, sklearn, pytest)",
     "  2. Exécution : pytest -q tests/",
     "",
     "Tests couverts :",
     "  • Cohérence des classes (3 classes)",
     "  • Distribution train/val/test",
     "  • Sparsité Lasso (L1)",
     "  • Accuracy baseline > seuil"],
    title_color=GREEN_DARK, font_size=13)

add_bullet_box(s, 6.6, 0.8, 6.2, 3.5,
    "Reproductibilité",
    ["requirements.txt complet et versionné",
     "  numpy>=1.26, tensorflow>=2.15",
     "  scikit-learn>=1.4, streamlit>=1.30",
     "  datasets>=2.19 (HuggingFace)",
     "",
     "Installation en 2 commandes :",
     "  python -m venv .venv",
     "  pip install -r requirements.txt",
     "",
     "Dataset téléchargé automatiquement",
     "  → aucun fichier de données versionné"],
    title_color=GREEN_DARK, font_size=13)

add_bullet_box(s, 0.5, 4.5, 12.3, 2.7,
    "Git Tags — Jalons tracés",
    ["data  →  Chargement dataset (commit initial)",
     "eda   →  Exploration & augmentation",
     "ml    →  Baseline ML + régularisation",
     "eval-ml  →  Évaluation ML + biais/variance",
     "dl       →  Architecture CNN + Transfer",
     "opti-dl  →  Optimisation (callbacks, Adam, BN)",
     "eval-dl  →  Pipeline complet + comparaison"],
    title_color=GREEN_DARK, font_size=13)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, GREEN_DARK)
rect(s, 0, 0, 13.33, 0.08, ACCENT)
rect(s, 0, 7.42, 13.33, 0.08, ACCENT)

box(s, 0.5, 0.5, 12.3, 0.7, "Conclusion", font_size=36, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
rect(s, 3.5, 1.3, 6.3, 0.05, ACCENT)

conclusions = [
    ("🎯 Objectif atteint",
     "Le Transfer Learning MobileNetV2 atteint 86.7% d'accuracy\nsoit +23 points vs la baseline ML"),
    ("📐 Modularité",
     "Code entièrement modulaire (src/*.py)\nimporté dans le notebook d'orchestration"),
    ("🚀 Déploiement",
     "Dashboard Streamlit fonctionnel\nprédiction en temps réel avec score de confiance"),
    ("🔬 Rigueur scientifique",
     "Analyse biais/variance, anti-vanishing gradient\ncomparaison argumentée ML vs DL"),
]
for i, (title, desc) in enumerate(conclusions):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.6 + row * 2.6
    rect(s, x, y, 6.0, 2.3, RGBColor(0x1B, 0x5E, 0x20))
    box(s, x+0.2, y+0.15, 5.7, 0.55, title, font_size=16, bold=True, color=ACCENT)
    box(s, x+0.2, y+0.75, 5.7, 1.3, desc, font_size=14, color=GREEN_LIGHT)

box(s, 0.5, 7.0, 12.3, 0.35,
    "MVOGO Abraham  ·  MAAROUFI Abdelhamid  —  PY-ML-DL-M1",
    font_size=13, color=GREEN_LIGHT, align=PP_ALIGN.CENTER, italic=True)


# ── Sauvegarde ────────────────────────────────────────────────────────────────
out = r"c:\Users\abrah\OneDrive\Bureau\IPSSI\beans-project\presentation_beans.pptx"
prs.save(out)
print(f"Fichier sauvegardé : {out}")
